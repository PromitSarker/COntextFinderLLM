import io
import csv
import json
from pathlib import Path
from typing import List, Dict, Any
from app.core.config import logger


class DocumentProcessor:
    """Processes non-PDF/non-image document types into text chunks."""

    def process_text_file(self, content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Process plain text, markdown, CSV, JSON, XML, HTML files."""
        try:
            ext = Path(filename).suffix.lower()
            text = ""

            if ext == ".csv":
                text = self._process_csv(content)
            elif ext == ".json":
                text = self._process_json(content)
            else:
                # .txt, .md, .html, .htm, .xml, .rtf
                text = content.decode("utf-8", errors="ignore")

            if not text.strip():
                logger.warning(f"No text extracted from {filename}")
                return []

            chunks = self._split_text(text, chunk_size=1000, overlap=200)

            documents = []
            for i, chunk in enumerate(chunks):
                documents.append({
                    "content": chunk,
                    "metadata": {
                        "filename": filename,
                        "source": "",
                        "page_number": i + 1,
                        "chunk_index": i,
                        "file_type": ext,
                    }
                })

            logger.info(f"Processed text file {filename}: {len(documents)} chunks")
            return documents

        except Exception as e:
            logger.error(f"Failed to process text file {filename}: {str(e)}")
            return []

    def process_docx(self, content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Process DOCX/DOC files."""
        try:
            import docx

            doc = docx.Document(io.BytesIO(content))
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        full_text.append(" | ".join(row_text))

            text = "\n".join(full_text)
            if not text.strip():
                logger.warning(f"No text extracted from {filename}")
                return []

            chunks = self._split_text(text, chunk_size=1000, overlap=200)

            documents = []
            for i, chunk in enumerate(chunks):
                documents.append({
                    "content": chunk,
                    "metadata": {
                        "filename": filename,
                        "source": "",
                        "page_number": i + 1,
                        "chunk_index": i,
                        "file_type": Path(filename).suffix.lower(),
                    }
                })

            logger.info(f"Processed DOCX {filename}: {len(documents)} chunks")
            return documents

        except ImportError:
            logger.error("python-docx is not installed. Run: pip install python-docx")
            return []
        except Exception as e:
            logger.error(f"Failed to process DOCX {filename}: {str(e)}")
            return []

    def process_pptx(self, content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Process PPTX/PPT files."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            documents = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_text.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                slide_text.append(" | ".join(row_text))

                if slide_text:
                    content_text = "\n".join(slide_text)
                    documents.append({
                        "content": content_text,
                        "metadata": {
                            "filename": filename,
                            "source": "",
                            "page_number": slide_num,
                            "chunk_index": slide_num - 1,
                            "file_type": Path(filename).suffix.lower(),
                        }
                    })

            logger.info(f"Processed PPTX {filename}: {len(documents)} slides")
            return documents

        except ImportError:
            logger.error("python-pptx is not installed. Run: pip install python-pptx")
            return []
        except Exception as e:
            logger.error(f"Failed to process PPTX {filename}: {str(e)}")
            return []

    def process_xlsx(self, content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Process XLSX/XLS files."""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            documents = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows_text = []

                for row in ws.iter_rows(values_only=True):
                    cell_values = [str(cell) for cell in row if cell is not None]
                    if cell_values:
                        rows_text.append(" | ".join(cell_values))

                if rows_text:
                    text = "\n".join(rows_text)
                    chunks = self._split_text(text, chunk_size=1000, overlap=200)

                    for i, chunk in enumerate(chunks):
                        documents.append({
                            "content": chunk,
                            "metadata": {
                                "filename": filename,
                                "source": "",
                                "page_number": i + 1,
                                "chunk_index": i,
                                "sheet_name": sheet_name,
                                "file_type": Path(filename).suffix.lower(),
                            }
                        })

            logger.info(f"Processed XLSX {filename}: {len(documents)} chunks")
            return documents

        except ImportError:
            logger.error("openpyxl is not installed. Run: pip install openpyxl")
            return []
        except Exception as e:
            logger.error(f"Failed to process XLSX {filename}: {str(e)}")
            return []

    def _process_csv(self, content: bytes) -> str:
        """Convert CSV content to readable text."""
        text_lines = []
        reader = csv.reader(io.StringIO(content.decode("utf-8", errors="ignore")))
        for row in reader:
            if any(cell.strip() for cell in row):
                text_lines.append(" | ".join(cell.strip() for cell in row))
        return "\n".join(text_lines)

    def _process_json(self, content: bytes) -> str:
        """Convert JSON content to readable text."""
        try:
            data = json.loads(content.decode("utf-8", errors="ignore"))
            return json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            return content.decode("utf-8", errors="ignore")

    def _split_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks."""
        if len(text) <= chunk_size:
            return [text]

        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size

            # Try to break at a sentence or paragraph boundary
            if end < len(text):
                # Look for paragraph break
                para_break = text.rfind("\n\n", start, end)
                if para_break > start + chunk_size // 2:
                    end = para_break + 2
                else:
                    # Look for sentence break
                    sentence_break = text.rfind(". ", start, end)
                    if sentence_break > start + chunk_size // 2:
                        end = sentence_break + 2

            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = end - overlap

        return chunks